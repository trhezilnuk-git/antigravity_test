import os
import sys
import collections.abc
from pptx import Presentation
import http.server
import socketserver
import threading
import webbrowser

HTML_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <title>PPTX Web Viewer</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
    <!-- Fonts and CSS libraries -->
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;800&display=swap" rel="stylesheet">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    
    <style>
        body { 
            font-family: 'Outfit', sans-serif; 
            background: #0f172a; /* Slate 900 */
            color: #f8fafc;
            overflow: hidden;
        }
        
        /* Modern Gradient Background Animation */
        .bg-mesh {
            position: absolute;
            top: 0; left: 0; width: 100vw; height: 100vh;
            z-index: -1;
            background: 
                radial-gradient(circle at 15% 50%, rgba(56, 189, 248, 0.15), transparent 50%),
                radial-gradient(circle at 85% 30%, rgba(168, 85, 247, 0.15), transparent 50%);
            animation: pulse-bg 15s ease-in-out infinite alternate;
        }
        
        @keyframes pulse-bg {
            0% { transform: scale(1); }
            100% { transform: scale(1.1); }
        }

        .reveal h2 { 
            font-family: 'Outfit', sans-serif;
            font-weight: 800;
            color: #38bdf8; /* Light Blue */
            text-transform: none;
            letter-spacing: -1px;
            text-shadow: 0 4px 15px rgba(56, 189, 248, 0.3);
            margin-bottom: 40px;
        }
        
        /* Glassmorphism Slide Content */
        .content-block { 
            margin-top: 20px; 
            text-align: center; 
            background: rgba(255, 255, 255, 0.03);
            padding: 50px 40px;
            border-radius: 24px;
            box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.3);
            backdrop-filter: blur(16px);
            -webkit-backdrop-filter: blur(16px);
            border: 1px solid rgba(255, 255, 255, 0.1);
            transition: all 0.3s ease;
        }
        
        .content-block:hover {
            background: rgba(255, 255, 255, 0.05);
            transform: translateY(-5px);
            border-color: rgba(56, 189, 248, 0.3);
        }

        .content-item { 
            margin-bottom: 20px; 
            font-size: 34px; 
            color: #e2e8f0; 
            line-height: 1.6; 
            font-weight: 300;
        }
        
        .content-item strong { color: #fff; font-weight: 600; }
        
        /* Reveal UI Elements */
        .reveal .controls { color: #a855f7; }
        .reveal .progress { color: #38bdf8; height: 4px; }
        .reveal .slide-number { 
            font-family: 'Outfit', sans-serif; 
            font-size: 20px !important; 
            color: rgba(255,255,255,0.6) !important; 
            background: transparent !important; 
        }
    </style>
</head>
<body>
    <div class="bg-mesh"></div>
    <div class="reveal">
        <div class="slides">
            {slides_html}
        </div>
    </div>
    
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            slideNumber: 'c/t',
            transition: 'convex',
            backgroundTransition: 'fade',
            center: true,
            controlsTutorial: true
        });
    </script>
</body>
</html>"""

def extract_pptx_data(file_path):
    print(f"Reading {file_path}...")
    prs = Presentation(file_path)
    slides_html = ""
    for slide in prs.slides:
        title = ""
        content = []
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text = shape.text.strip()
            if not text: continue
            if shape == slide.shapes.title:
                title = text
            else:
                content.append(text)
                
        slide_html = "<section>"
        if title:
            slide_html += f"<h2>{title}</h2>"
        if content:
            slide_html += '<div class="content-block">'
            for c in content:
                c = c.replace('\\n', '<br>')
                slide_html += f'<p class="content-item">{c}</p>'
            slide_html += '</div>'
        slide_html += "</section>\\n"
        slides_html += slide_html
    return slides_html

def main():
    pptx_file = "test_10_slides.pptx"
    if not os.path.exists(pptx_file):
        print(f"Error: {pptx_file} not found. Please ensure it was created.")
        sys.exit(1)

    print("Extracting slides data from PPTX...")
    slides_html = extract_pptx_data(pptx_file)
    html_content = HTML_TEMPLATE.replace("{slides_html}", slides_html)
    
    with open("index.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    
    print("Generated index.html successfully.")
    
    PORT = 8080
    
    class CustomHandler(http.server.SimpleHTTPRequestHandler):
        def log_message(self, format, *args):
            pass # Suppress log output for a cleaner console
            
    with socketserver.TCPServer(("", PORT), CustomHandler) as httpd:
        url = f"http://localhost:{PORT}"
        print(f"\\nWebsite serves at {url}")
        print("Opening browser automatically...")
        print("Press Ctrl+C to stop the server.\\n")
        
        # Open in browser automatically
        threading.Timer(1.0, lambda: webbrowser.open(url)).start()
        
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print("\\nShutting down server...")
            httpd.shutdown()

if __name__ == "__main__":
    main()
