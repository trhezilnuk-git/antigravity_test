from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Top 10 Modern AI Tools"
subtitle.text = "An Overview of the Best AI Capabilities"

# 10 tools
ai_tools = [
    {
        "name": "1. ChatGPT (OpenAI)",
        "content": "Description: The most popular conversational AI.\\nKey Features:\\n- Advanced natural language understanding\\n- Code generation and debugging\\n- Data analysis and visualization"
    },
    {
        "name": "2. Claude 3 (Anthropic)",
        "content": "Description: A highly capable AI focusing on safety and large context.\\nKey Features:\\n- Massive context window (200k+ tokens)\\n- Excellent at reading long documents\\n- Nuanced and safer responses"
    },
    {
        "name": "3. Gemini (Google)",
        "content": "Description: Google's flagship multimodal AI model.\\nKey Features:\\n- Native multimodal support (text, image, audio, video)\\n- Deep integration with Google Workspace\\n- Real-time information access"
    },
    {
        "name": "4. Midjourney",
        "content": "Description: Leading AI image generation platform.\\nKey Features:\\n- Unbelievably high-quality images\\n- Incredibly artistic interpretations\\n- Operated primarily via Discord"
    },
    {
        "name": "5. GitHub Copilot",
        "content": "Description: The most popular AI pair programmer.\\nKey Features:\\n- Seamless IDE integration\\n- Real-time autocomplete for code\\n- Translates natural language to code"
    },
    {
        "name": "6. Perplexity AI",
        "content": "Description: AI-powered research and search engine.\\nKey Features:\\n- Real-time web citations\\n- Focuses on accuracy and sourcing\\n- Great for deep-dive research"
    },
    {
        "name": "7. Sora (OpenAI)",
        "content": "Description: Groundbreaking AI video generation model.\\nKey Features:\\n- Generates up to 1-minute realistic videos\\n- Follows complex prompt instructions\\n- Maintains physical coherence"
    },
    {
        "name": "8. Suno AI",
        "content": "Description: Advanced AI music generation tool.\\nKey Features:\\n- Generates full songs with vocals and instruments\\n- Variety of genres and styles\\n- High audio quality"
    },
    {
        "name": "9. Notion AI",
        "content": "Description: AI seamlessly integrated into Notion.\\nKey Features:\\n- Brainstorming and drafting\\n- Summarizing existing notes\\n- Translating and organizing content"
    },
    {
        "name": "10. ElevenLabs",
        "content": "Description: State-of-the-art AI voice generator.\\nKey Features:\\n- Extremely realistic text-to-speech\\n- Voice cloning capabilities\\n- Emotional range and inflection"
    }
]

for tool in ai_tools:
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = tool["name"]
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = tool["content"]

prs.save('test_10_slides.pptx')
print('Successfully created test_10_slides.pptx')
